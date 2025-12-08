"""
File Text Extraction Service
Handles text extraction from various file formats
"""

import os
from typing import Optional
from io import BytesIO

# PDF
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

# DOCX
try:
    import docx
except ImportError:
    docx = None

# CSV/XLSX
try:
    import pandas as pd
except ImportError:
    pd = None

# PPTX
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Images (OCR)
try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None

# PDF to image for OCR
try:
    from pdf2image import convert_from_path
except ImportError:
    convert_from_path = None

class TextExtractor:
    """Service for extracting text from various file formats"""
    
    def __init__(self):
        """Initialize text extractor"""
        self.supported_formats = {
            'pdf': self._extract_pdf_text,
            'docx': self._extract_docx_text,
            'txt': self._extract_txt_text,
            'csv': self._extract_csv_text,
            'xlsx': self._extract_xlsx_text,
            'pptx': self._extract_pptx_text
        }
    
    def extract_text_from_file(self, file_path: str) -> str:
        """Extract text from a file on disk"""
        ext = os.path.splitext(file_path)[-1].lower()
        print(f"[DEBUG] Extracting text from: {file_path} (ext: {ext})")
        
        try:
            if ext == ".pdf" and PyPDF2:
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    text = "\n".join(page.extract_text() or "" for page in reader.pages)
                if text.strip():
                    return text
                # If no text found, try OCR
                if convert_from_path and pytesseract:
                    try:
                        images = convert_from_path(file_path)
                        ocr_text = []
                        for img in images:
                            ocr_text.append(pytesseract.image_to_string(img))
                        return "\n".join(ocr_text) or "Text extraction not supported for this file type, but file is stored."
                    except Exception as e:
                        print(f"[DEBUG] Error extracting text with OCR: {e}")
                        return f"Error extracting text with OCR: {e}"
                else:
                    return "Text extraction not supported for this file type, but file is stored."
            elif ext == ".docx" and docx:
                doc = docx.Document(file_path)
                return "\n".join(p.text for p in doc.paragraphs)
            elif ext == ".txt":
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()
            elif ext == ".csv" and pd:
                df = pd.read_csv(file_path)
                return df.to_string()
            elif ext == ".xlsx" and pd:
                df = pd.read_excel(file_path)
                return df.to_string()
            elif ext == ".pptx" and Presentation:
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)
            elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"] and Image and pytesseract:
                img = Image.open(file_path)
                return pytesseract.image_to_string(img)
            else:
                print(f"[DEBUG] Unsupported file type or missing dependency for: {file_path}")
                return "Text extraction not supported for this file type, but file is stored."
        except Exception as e:
            print(f"[DEBUG] Error extracting text: {e}")
            return f"Error extracting text: {e}"
    
    def extract_text_from_bytes(self, file_bytes: BytesIO, file_type: str) -> str:
        """Extract text from an in-memory file (BytesIO) given its type"""
        try:
            ext = file_type.lower()
            if not ext.startswith('.'):
                ext = '.' + ext
            
            if ext == ".pdf" and PyPDF2:
                reader = PyPDF2.PdfReader(file_bytes)
                text = "\n".join(page.extract_text() or "" for page in reader.pages)
                return text
            elif ext == ".docx" and docx:
                doc = docx.Document(file_bytes)
                return "\n".join(p.text for p in doc.paragraphs)
            elif ext == ".txt":
                return file_bytes.read().decode("utf-8", errors="ignore")
            elif ext == ".csv" and pd:
                return pd.read_csv(file_bytes).to_string()
            elif ext == ".xlsx" and pd:
                return pd.read_excel(file_bytes).to_string()
            elif ext == ".pptx" and Presentation:
                prs = Presentation(file_bytes)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)
            elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"] and Image and pytesseract:
                img = Image.open(file_bytes)
                return pytesseract.image_to_string(img)
            else:
                return "Text extraction not supported for this file type."
        except Exception as e:
            return f"Error extracting text: {e}"
    
    def _extract_pdf_text(self, file_bytes: BytesIO) -> str:
        """Extract text from PDF"""
        try:
            pdf_reader = PyPDF2.PdfReader(file_bytes)
            text_parts = []
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                except Exception as e:
                    print(f"Warning: Error extracting text from PDF page {page_num}: {str(e)}")
                    text_parts.append(f"[Error on page {page_num + 1}]")
            
            return "\n".join(text_parts) if text_parts else "[No text extracted from PDF]"
        
        except Exception as e:
            print(f"PDF extraction failed: {str(e)}")
            return f"[PDF extraction error: {str(e)}]"
    
    def _extract_docx_text(self, file_bytes: BytesIO) -> str:
        """Extract text from DOCX"""
        try:
            doc = docx.Document(file_bytes)
            text_parts = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            return "\n".join(text_parts) if text_parts else "[No text extracted from DOCX]"
        
        except Exception as e:
            print(f"DOCX extraction failed: {str(e)}")
            return f"[DOCX extraction error: {str(e)}]"
    
    def _extract_txt_text(self, file_bytes: BytesIO) -> str:
        """Extract text from TXT"""
        try:
            content = file_bytes.read().decode('utf-8', errors='ignore')
            return content if content.strip() else "[Empty text file]"
        except Exception as e:
            print(f"TXT extraction failed: {str(e)}")
            return f"[TXT extraction error: {str(e)}]"
    
    def _extract_xlsx_text(self, file_bytes: BytesIO) -> str:
        """Extract text from XLSX"""
        try:
            df = pd.read_excel(file_bytes, engine='openpyxl')
            return df.to_string() if not df.empty else "[Empty Excel file]"
        except Exception as e:
            print(f"XLSX extraction failed: {str(e)}")
            return f"[XLSX extraction error: {str(e)}]"
    
    def _extract_csv_text(self, file_bytes: BytesIO) -> str:
        """Extract text from CSV"""
        try:
            df = pd.read_csv(file_bytes)
            return df.to_string() if not df.empty else "[Empty CSV file]"
        except Exception as e:
            print(f"CSV extraction failed: {str(e)}")
            return f"[CSV extraction error: {str(e)}]"
    
    def _extract_pptx_text(self, file_bytes: BytesIO) -> str:
        """Extract text from PPTX"""
        try:
            prs = Presentation(file_bytes)
            text_parts = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            
            return "\n".join(text_parts) if text_parts else "[No text extracted from PPTX]"
        
        except Exception as e:
            print(f"PPTX extraction failed: {str(e)}")
            return f"[PPTX extraction error: {str(e)}]"
    
    def get_supported_formats(self) -> list:
        """Get list of supported file formats"""
        return list(self.supported_formats.keys())
    
    def is_format_supported(self, file_type: str) -> bool:
        """Check if a file format is supported"""
        return file_type.lower() in self.supported_formats 